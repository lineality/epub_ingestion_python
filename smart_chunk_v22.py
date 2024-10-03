###################################################
# This code automatically finds and processes
# epub books, docx, pdf, and plain text docs
# especially for RAG document ingestion processing
###################################################
"""
Instructions:
(Optional set/change configuration below.)
1. put target files (or directories of files) in a folder called /target_files/ with this program
2. run this program
3. Answer attribution text question for each file
4. results in /txt_pool/ directory


input -> one or more epub/pdf/docx/pptx/txt files
output -> txt and json files that contain the text from sections of the book
          as well as smart-chunked segments made to your size specs
          e.g. a max of 500 characters, which contain whole sentences.
          Chunks do not cut words or sentences in half.

# Set of Results, saved in a file per epub doc:
1. One .jsonl file
2. (Plural) Individual .json files in a folder
3. One .txt text file of the whole epub
4. (plural) Individual .txt files from separate parts of epub
5. chunks as one .jsonl file
6. (Plural) chunks under specified character length as separate text
Future Feature:
7. Chunk-Metadata (for model training, for DB-retrieval, etc.)

There is the option of removing small files, but appending may often be a better default. 

You may need to adjust MAX_CHUNK_SIZE (biggest file size) and MINIMUM_BYTES_SIZE (smallest file size)
over a few tries to make sure the results fit in your workflow.
"""
################
# Configuration
################
# size character length of chunk
MAX_CHUNK_SIZE = 1800
MINIMUM_BYTES_SIZE = 300
REMOVE_BELOW_SIZE = (
    0  # Default is zero, don't remove anything (nothing is smaller than zero)
)
RESULTS_DIR_NAME = "ingestion_processing_results"

# if needed, set PDF-reader below, 'all' is default
PDF_USE_ALL = False
PDF_TRY_PYMU = True
PDF_TRY_PYPDF = False
PDF_TRY_PDFPLUMBER = False

pool_counter = 1

import zipfile
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import json
import os
import shutil
import re
import time
import traceback
import docx
from pypdf import PdfReader
import pdfplumber
import fitz
import glob
from datetime import datetime
from pptx import Presentation

"""
requirements.txt ->
python-pptx
python-docx
PyMuPDF
pypdf
pdfplumber
beautifulsoup4

pip install PyMuPDF
pip install pypdf
pip install pdfplumber
pip install python-pptx
pip install python-docx
pip install beautifulsoup4

python -m pip install python-pptx python-docx PyMuPDF pypdf pdfplumber

fitz = https://pypi.org/project/PyMuPDF/
https://pypi.org/project/pypdf/
https://pypi.org/project/pdfplumber/
https://pypi.org/project/python-pptx/
"""


# start debug timer
start_time_outer = time.monotonic()
start_time_inner = time.time()


def remove_small_files(directory, size):
    """
    Remove files from a specified directory that are smaller than a given size.

    Parameters:
    directory (str): The path to the directory.
    size (int): The size threshold in bytes. Files smaller than this size will be removed.

    Returns:
    None

    Requires:
        import os
    """
    counter = 0

    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        if os.path.isfile(file_path) and os.path.getsize(file_path) < size:
            os.remove(file_path)
            counter += 1

    print(f"Removed {counter} files from The Mines of Moria .../{directory}")


def count_files(directory):
    """
    count_files from a specified directory

    Parameters:
    directory (str): The path to the directory.

    Returns:
    None

    Requires:
        import os
    """
    counter = 0

    for filename in os.listdir(directory):

        counter += 1

    print(
        f"Current File Population in the Land of {directory} observed to be at: {counter}"
    )


# def get_files_from_subdirs(base_dir, file_extension):
#     """
#     Get the paths of all files with the specified extension from subdirectories in the base directory.

#     Args:
#         base_dir (str): The path to the base directory containing the subdirectories.
#         file_extension (str): The file extension to match (e.g., '.txt', '.py', '.jpg').

#     Returns:
#         list: A list of file paths matching the specified extension.
#     """
#     # Get all the subdirectories in the base directory
#     subdirs = [d for d in os.listdir(base_dir) if os.path.isdir(os.path.join(base_dir, d))]

#     # Get the paths of all files with the specified extension in the subdirectories
#     file_paths = []
#     for subdir in subdirs:
#         subdir_path = os.path.join(base_dir, subdir)
#         file_paths.extend(glob.glob(os.path.join(subdir_path, f'*{file_extension}')))

#     return file_paths

# def get_files_from_subdirs(base_dir, file_extension):
#     """
#     Get the paths of all files with the specified extension from subdirectories and the base directory.

#     Args:
#         base_dir (str): The path to the base directory containing the subdirectories.
#         file_extension (str): The file extension to match (e.g., '.txt', '.py', '.jpg').

#     Returns:
#         list: A list of file paths matching the specified extension.

#     Requires:
#         - import os
#         - import glob

#     """
#     # Ensure the file extension starts with a dot
#     if not file_extension.startswith('.'):
#         file_extension = '.' + file_extension

#     # Get all the subdirectories in the base directory
#     subdirs = [d for d in os.listdir(base_dir) if os.path.isdir(os.path.join(base_dir, d))]

#     # Get the paths of all files with the specified extension in the subdirectories
#     file_paths = []
#     for subdir in subdirs:
#         subdir_path = os.path.join(base_dir, subdir)
#         file_paths.extend(glob.glob(os.path.join(subdir_path, f'*{file_extension}')))

#     # Get the paths of all files with the specified extension in the base directory
#     base_dir_files = glob.glob(os.path.join(base_dir, f'*{file_extension}'))
#     file_paths.extend(base_dir_files)

#     return file_paths


def get_files_from_subdirs(base_dir, file_extension):
    """
    Get the paths of all files with the specified extension from subdirectories and the base directory.

    Args:
        base_dir (str): The path to the base directory containing the subdirectories.
        file_extension (str): The file extension to match (e.g., '.txt', '.py', '.jpg').

    Returns:
        list: A list of file paths matching the specified extension.

    Requires:
        - import os
        - import glob
    """
    # Ensure the file extension starts with a dot
    if not file_extension.startswith("."):
        file_extension = "." + file_extension

    # Get the paths of all files with the specified extension in all subdirectories (sub-sub, etc)
    file_paths = []
    for root, dirs, files in os.walk(base_dir):
        file_paths.extend(glob.glob(os.path.join(root, f"*{file_extension}")))

    # Get the paths of all files with the specified extension in the base directory
    base_dir_files = glob.glob(os.path.join(base_dir, f"*{file_extension}"))
    file_paths.extend(base_dir_files)

    return file_paths


def remove_control_chars(s):
    # Create a translation table that maps all control characters to None
    control_chars = dict.fromkeys(range(0, 32), None)
    # Add the U+000B character explicitly
    control_chars[11] = None
    # Translate the string using the translation table
    return s.translate(control_chars)


def extract_pptx_text_to_file(pptx_file):
    """
    extract text from a .pptx slide deck
    ideally,
    extract the text from each slide,
    """
    # Get the PowerPoint file name without extension
    file_name = os.path.splitext(os.path.basename(pptx_file))[0]

    # Specify the directory path you want to create
    save_here_directory_path = f"pptx_chunks/{file_name}_chunks"

    # Check if the directory exists
    if not os.path.exists(save_here_directory_path):
        # If the directory doesn't exist, create it
        os.makedirs(save_here_directory_path)
        print(f"Directory '{save_here_directory_path}' created successfully.")
    else:
        print(f"Directory '{save_here_directory_path}' already exists.")

    # Open the PowerPoint file
    prs = Presentation(pptx_file)

    # Get the current timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    slide_number = 1

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # f.write(shape.text + "\n")

                raw_text = shape.text
                utf8_text = remove_control_chars(raw_text)

                # save slide text
                slide_text.append(utf8_text)

        # if not empty
        if slide_text:

            # # Inspection
            # print("slide_text")
            # print(slide_text)

            # Create the output file name with slide number and timestamp
            output_file = (
                f"{save_here_directory_path}/{file_name}_{slide_number}_{timestamp}.txt"
            )

            # Write slide text to the output file
            with open(output_file, "w", encoding="utf-8") as f:
                f.write("\n".join(slide_text))

            print(f"Slide {slide_number} text saved to {output_file}")

        slide_number += 1

    print(f"Text extracted and saved to {output_file}")


def simple_extracttextfrom_pdf(pdf_path):
    """
    Try three methods for reading pdf,
    returns longest, in case of partial-fails

    # Example usage
    pdf_path = '2019 PMR Inn Type.pdf'
    extracted_text = extract_text_from_pdf(pdf_path)
    print(extracted_text)
    """

    # in case none are created
    text_pypdf = ""
    text_pdfplumber = ""
    text_pymupdf = ""

    if PDF_TRY_PYMU or PDF_USE_ALL:
        try:
            # PyMuPDF
            print("Trying PyMuPDF...")
            pdf_file = fitz.open(pdf_path)
            text_pymupdf = ""
            for page_num in range(len(pdf_file)):
                page = pdf_file[page_num]
                text_pymupdf += page.get_text()
        except fitz.FileDataError:
            text_pymupdf = ""
        except Exception as e:
            print(f"Error (PyMuPDF): {e}")
            text_pymupdf = ""

    if PDF_TRY_PYPDF or PDF_USE_ALL:
        try:
            # pypdf
            print("Trying pypdf PdfReader...")
            reader = PdfReader(pdf_path)
            text_pypdf = ""
            for page in reader.pages:
                text_pypdf += page.extract_text()
        except Exception as e:
            print(f"Error (pypdf): {e}")
            text_pypdf = ""

    if PDF_TRY_PDFPLUMBER or PDF_USE_ALL:
        try:
            # pdfplumber
            print("Trying pdfplumber...")
            with pdfplumber.open(pdf_path) as pdf:
                text_pdfplumber = ""
                for page in pdf.pages:
                    text_pdfplumber += page.extract_text()
        except Exception as e:
            print(f"Error (pdfplumber): {e}")
            text_pdfplumber = ""

    # Keep the longest extracted text
    print(
        f""" len()
    text_pypdf      -> {len(text_pypdf)}
    text_pdfplumber -> {len(text_pdfplumber)}
    text_pymupdf    -> {len(text_pymupdf)}
    """
    )

    # Keep the longest extracted text
    longest_text = max([text_pypdf, text_pdfplumber, text_pymupdf], key=len)

    # Check if any text was extracted
    if not longest_text:
        print("Error: Unable to extract text from the PDF file.")
        return None

    return longest_text


def get_ordered_html_files(opf_content):
    """
    Parses the content.opf file to determine the reading order of HTML files in the EPUB.

    The function reads the 'content.opf' file, which contains metadata about the EPUB's structure.
    It identifies the 'spine' element, which lists the reading order of the content documents,
    and the 'manifest' element, which provides the location of these documents.
    The function returns a list of HTML file paths in the order they should be read.

    Args:
    opf_content (str): A string representation of the content.opf file.

    Returns:
    list: An ordered list of HTML file paths as specified in the EPUB's spine.
    """

    # Parse the content.opf XML content
    tree = ET.ElementTree(ET.fromstring(opf_content))
    root = tree.getroot()

    # Define the namespace for the OPF package file
    ns = {"opf": "http://www.idpf.org/2007/opf"}

    # Find the spine element which indicates the order of the content documents
    spine = root.find("opf:spine", ns)
    itemrefs = spine.findall("opf:itemref", ns)

    # Extract the id references for each item in the spine
    item_ids = [itemref.get("idref") for itemref in itemrefs]

    # Find the manifest element which lists all the content documents
    manifest = root.find("opf:manifest", ns)
    items = manifest.findall("opf:item", ns)

    # Create a dictionary mapping item IDs to their corresponding file paths
    html_files = {
        item.get("id"): item.get("href")
        for item in items
        if item.get("media-type") == "application/xhtml+xml"
    }

    # Generate an ordered list of HTML files based on the spine order
    ordered_html_files = [
        html_files[item_id] for item_id in item_ids if item_id in html_files
    ]

    return ordered_html_files


def extract_text_from_html(html_content, this_epub_output_dir_path):
    """
    Extracts and returns text from an HTML content.
    """
    # print("HTML Content before BeautifulSoup Parsing:\n", html_content[:500])  # Print first 500 characters of HTML
    print_and_log(
        f"\nlen(HTML Content before BeautifulSoup Parsing) -> {len(html_content)}",
        this_epub_output_dir_path,
    )  # Print first 500 characters of HTML

    soup = BeautifulSoup(html_content, "html.parser")
    parsed_text = soup.get_text()
    # print("Extracted Text:\n", parsed_text[:500])  # Print first 500 characters of extracted text
    print_and_log(
        f"len(Extracted Text) -> {len(parsed_text)}", this_epub_output_dir_path
    )  # Print first 500 characters of extracted text

    return parsed_text


def fix_text_formatting(text):
    """Replaces the Unicode right single quotation mark with a standard apostrophe."""
    return text.replace("\u2019", "'")


def check_len_chunks_in_list(chunks_list, max_chunk_size, this_epub_output_dir_path):

    size_flag_ok = True

    for index, this_chunk in enumerate(chunks_list):

        # get size of chunk
        this_length = len(this_chunk)

        # check size against user-input max size
        if this_length > max_chunk_size:
            print_and_log(this_length, this_epub_output_dir_path)
            print_and_log(
                f"""
            Warning: chunk over max size.
            This chunk size: {len(this_chunk)}.
            Max size: {max_chunk_size}
            """,
                this_epub_output_dir_path,
            )
            print_and_log(f"This chunk: {len(this_chunk)}", this_epub_output_dir_path)
            size_flag_ok = False

    # report and log
    if size_flag_ok:
        print_and_log("Size Check, OK \\o/", this_epub_output_dir_path)
    else:
        print_and_log("WARNING: Size Check Failed!", this_epub_output_dir_path)


def save_individual_chunks(chunks_list, output_chunks_dir, chunk_source_name):
    # Strip non-alphanumeric characters from the directory name
    # Remove spaces from the output_chunks_dir path
    new_output_chunks_dir = re.sub(r"\s+", "_", output_chunks_dir)

    # Get the absolute path of the output_chunks_dir
    new_output_chunks_dir = os.path.abspath(new_output_chunks_dir)

    # print(f"new_output_chunks_dir -> {new_output_chunks_dir}")

    if not os.path.exists(new_output_chunks_dir):
        try:
            os.makedirs(new_output_chunks_dir)
        except Exception as e:
            print_and_log(traceback.format_exc())  # This will print the stack trace
            raise e

    for index, this_chunk in enumerate(chunks_list):
        chunk_name = f"{chunk_source_name}_{index}.txt"

        # remove spaces
        chunk_name = re.sub(r"\s+", "_", chunk_name)

        # # Save individual txt files
        # print('new_output_chunks_dir -> ', new_output_chunks_dir)
        # print('chunk_name -> ', chunk_name)

        # Get just the file name
        chunk_name = os.path.basename(chunk_name)

        # print('chunk_name -> ', chunk_name)
        # print(f"os.path.exists(new_output_chunks_dir) -> {os.path.exists(new_output_chunks_dir)}")

        individual_chunk_path = os.path.join(new_output_chunks_dir, chunk_name)
        # print('individual_chunk_path -> ', individual_chunk_path)

        with open(individual_chunk_path, "w") as f:
            f.write(this_chunk)
            # print('chunk written', individual_chunk_path)

        ###################
        # save to txt.pool
        ###################

        pool_output_chunks_dir = "txt_pool"

        if not os.path.exists(pool_output_chunks_dir):
            try:
                os.makedirs(pool_output_chunks_dir)
            except Exception as e:
                print_and_log(traceback.format_exc())  # This will print the stack trace
                raise e

        individual_chunk_path = os.path.join(pool_output_chunks_dir, chunk_name)
        print("pool_output_chunks_dir -> ", pool_output_chunks_dir)
        print("individual_chunk_path -> ", individual_chunk_path)

        with open(individual_chunk_path, "w") as f:
            f.write(this_chunk)

    return len(chunks_list)


def append_chunks_to_jsonl(chunks_list, output_chunks_jsonl_path, chunk_source_name):
    """Appends chunks of text to a .jsonl file, each chunk as a JSON object.

    Args:
        chunks_list (list): List of text chunks to be appended.
        output_jsonl_path (str): The output file path for the .jsonl file.
        chunk_source_name (str): Base name for each chunk, used in the 'source_name' field.

    Returns:
        int: The number of chunks appended.
    """

    with open(output_chunks_jsonl_path, "a") as f:  # Open file in append mode
        for index, this_chunk in enumerate(chunks_list):
            # Construct a JSON object for the chunk
            chunk_data = {
                "source_name": f"{chunk_source_name}_{index}",
                "text": this_chunk,
            }

            # Convert the chunk data to a JSON string and append it to the file with a newline
            f.write(json.dumps(chunk_data) + "\n")

    return len(chunks_list)


def print_and_log(input_text, this_epub_output_dir_path):
    # check if input is a string, if not...make it a string!
    if not isinstance(input_text, str):
        input_text = str(input_text)

    # print to terminal
    print(input_text)

    # log file path is...
    log_file_path = os.path.join(this_epub_output_dir_path, "log.txt")

    # log: Write/Append to a log.txt file
    with open(log_file_path, "a") as f:
        f.write(input_text + "\n\n")


def extract_text_from_epub(
    epub_file_path,
    this_txt_output_dir_path,
    output_jsonl_path,
    output_json_dir,
    output_whole_txt_path,
    output_txt_dir,
    output_chunks_jsonl_path,
    output_chunks_dir,
    max_chunk_size=MAX_CHUNK_SIZE,
):

    with zipfile.ZipFile(epub_file_path, "r") as epub:
        print_and_log(f"EPUB Contents: -> {epub.namelist()}", this_epub_output_dir_path)

        ###################
        # Make Directories
        ###################

        # Create a directory for individual JSON files
        if not os.path.exists(output_json_dir):
            os.makedirs(output_json_dir)

        # Create a directory for individual txt files
        if not os.path.exists(output_txt_dir):
            os.makedirs(output_txt_dir)

        # Create a directory for chunks output_chunks_dir
        if not os.path.exists(output_chunks_dir):
            os.makedirs(output_chunks_dir)

        ##################################
        # Get & Read html files from epub
        ##################################
        # find opf file
        opf_file = [f for f in epub.namelist() if "content.opf" in f][0]

        # read opf file
        opf_content = epub.read(opf_file).decode("utf-8")

        # get ordered HTML files
        ordered_html_files_list = get_ordered_html_files(opf_content)

        ############################################
        # Read and extract text from each HTML file
        ############################################

        # iterate through html files
        for html_file in ordered_html_files_list:
            full_path = os.path.join(os.path.dirname(opf_file), html_file)
            if full_path in epub.namelist():
                html_content = epub.read(full_path).decode("utf-8")

                #########################
                # extract text from epub
                #########################
                raw_text = extract_text_from_html(
                    html_content, this_epub_output_dir_path
                )
                print_and_log(
                    f"len(text for json)-> {len(raw_text)}", this_epub_output_dir_path
                )

                # fix text formatting
                text = fix_text_formatting(raw_text)

                #################
                # .json & .jsonl
                #################

                # Write/Append to a single JSONL file
                with open(output_jsonl_path, "a") as f:
                    json_record = json.dumps({"text": text.strip()})
                    f.write(json_record + "\n")

                # Save individual JSON file
                individual_json_path = os.path.join(
                    output_json_dir, f"{os.path.splitext(html_file)[0]}.json"
                )
                with open(individual_json_path, "w") as f:
                    json.dump({"text": text.strip()}, f, indent=4)

                #######
                # .txt
                #######

                # Write/Append to a single text .txt file
                with open(output_whole_txt_path, "a") as f:
                    f.write(text + "\n\n")

                # Save individual txt files
                individual_txt_path = os.path.join(
                    output_txt_dir, f"{os.path.splitext(html_file)[0]}.txt"
                )
                with open(individual_txt_path, "w") as f:
                    f.write(text)

                #########
                # chunks
                #########

                chunks_list = make_chunk_list(
                    text, max_chunk_size, this_epub_output_dir_path
                )

                chunk_source_name = os.path.splitext(html_file)[0]

                # check sizes
                check_len_chunks_in_list(
                    chunks_list, max_chunk_size, this_epub_output_dir_path
                )

                number_of_chunks = save_individual_chunks(
                    chunks_list, output_chunks_dir, chunk_source_name
                )
                print_and_log(
                    f"Chunked: split into this many chunks-> {number_of_chunks}",
                    this_epub_output_dir_path,
                )

                append_chunks_to_jsonl(
                    chunks_list, output_chunks_jsonl_path, chunk_source_name
                )

                print_and_log(f"{html_file} -> ok!", this_epub_output_dir_path)

            else:  # File Not Found
                print_and_log(
                    f"Warning: File {full_path} not found in the archive.",
                    this_epub_output_dir_path,
                )


def extract_text_from_txt(
    text_file_path,
    this_txt_output_dir_path,
    output_jsonl_path,
    output_json_dir,
    output_whole_txt_path,
    output_txt_dir,
    output_chunks_jsonl_path,
    output_chunks_dir,
    max_chunk_size=MAX_CHUNK_SIZE,
):

    # Open the file in read mode
    with open(text_file_path, "r") as file:

        ###################
        # Make Directories
        ###################

        # Create a directory for individual JSON files
        if not os.path.exists(output_json_dir):
            os.makedirs(output_json_dir)

        # Create a directory for individual txt files
        if not os.path.exists(output_txt_dir):
            os.makedirs(output_txt_dir)

        # Create a directory for chunks output_chunks_dir
        if not os.path.exists(output_chunks_dir):
            os.makedirs(output_chunks_dir)

        #########################
        # extract text from txt
        #########################
        # Read the entire contents of the file
        text = file.read()

        if text:
            # Save individual txt files
            text_file_path = os.path.basename(text_file_path)

            #################
            # .json & .jsonl
            #################

            # Write/Append to a single JSONL file
            with open(output_jsonl_path, "a") as f:
                json_record = json.dumps({"text": text.strip()})
                f.write(json_record + "\n")

            # Save individual JSON file
            individual_json_path = os.path.join(
                output_json_dir, f"{os.path.splitext(text_file_path)[0]}.json"
            )
            with open(individual_json_path, "w") as f:
                json.dump({"text": text.strip()}, f, indent=4)

            #######
            # .txt
            #######

            # Write/Append to a single text .txt file
            with open(output_whole_txt_path, "a") as f:
                f.write(text + "\n\n")

            # Save individual txt files
            individual_txt_path = os.path.join(
                output_txt_dir, f"{os.path.splitext(text_file_path)[0]}.txt"
            )
            with open(individual_txt_path, "w") as f:
                f.write(text)

            #########
            # chunks
            #########

            chunks_list = make_chunk_list(
                text, max_chunk_size, this_txt_output_dir_path
            )

            chunk_source_name = os.path.splitext(text_file_path)[0]

            # check sizes
            check_len_chunks_in_list(
                chunks_list, max_chunk_size, this_txt_output_dir_path
            )

            number_of_chunks = save_individual_chunks(
                chunks_list, output_chunks_dir, chunk_source_name
            )
            print_and_log(
                f"Chunked: split into this many chunks-> {number_of_chunks}",
                this_txt_output_dir_path,
            )

            append_chunks_to_jsonl(
                chunks_list, output_chunks_jsonl_path, chunk_source_name
            )

            print_and_log(f"{text_file_path} -> ok!", this_txt_output_dir_path)

            print("OK!")

        else:
            print_and_log(
                f"{text_file_path} -> Faile, no text extracted",
                this_txt_output_dir_path,
            )


def extract_text_from_docx(
    text_file_path,
    this_txt_output_dir_path,
    output_jsonl_path,
    output_json_dir,
    output_whole_txt_path,
    output_txt_dir,
    output_chunks_jsonl_path,
    output_chunks_dir,
    max_chunk_size=MAX_CHUNK_SIZE,
):
    """
    Extracts the text from a Microsoft Word (DOCX) file.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        str: The extracted text from the DOCX file.
    """

    ###################
    # Make Directories
    ###################

    # Create a directory for individual JSON files
    if not os.path.exists(output_json_dir):
        os.makedirs(output_json_dir)

    # Create a directory for individual txt files
    if not os.path.exists(output_txt_dir):
        os.makedirs(output_txt_dir)

    # Create a directory for chunks output_chunks_dir
    if not os.path.exists(output_chunks_dir):
        os.makedirs(output_chunks_dir)

    #########################
    # extract text from docx
    #########################
    doc = docx.Document(text_file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    if text:
        # Save individual txt files
        text_file_path = os.path.basename(text_file_path)

        #################
        # .json & .jsonl
        #################

        # Write/Append to a single JSONL file
        with open(output_jsonl_path, "a") as f:
            json_record = json.dumps({"text": text.strip()})
            f.write(json_record + "\n")

        # Save individual JSON file
        individual_json_path = os.path.join(
            output_json_dir, f"{os.path.splitext(text_file_path)[0]}.json"
        )
        with open(individual_json_path, "w") as f:
            json.dump({"text": text.strip()}, f, indent=4)

        #######
        # .txt
        #######

        # Write/Append to a single text .txt file
        with open(output_whole_txt_path, "a") as f:
            f.write(text + "\n\n")

        # # Save individual txt files
        # individual_txt_path = os.path.join(output_txt_dir, f"{os.path.splitext(text_file_path)[0]}.txt")
        # print(f"individual_txt_path -> {individual_txt_path}")
        # with open(individual_txt_path, 'w') as f:
        #     f.write(text)

        individual_txt_path = os.path.join(output_txt_dir, text_file_path)
        print(f"individual_txt_path -> {individual_txt_path}")
        with open(individual_txt_path, "w", encoding="utf-8") as f:
            f.write(text)

        #########
        # chunks
        #########

        chunks_list = make_chunk_list(text, max_chunk_size, this_txt_output_dir_path)

        chunk_source_name = os.path.splitext(text_file_path)[0]

        # check sizes
        check_len_chunks_in_list(chunks_list, max_chunk_size, this_txt_output_dir_path)

        number_of_chunks = save_individual_chunks(
            chunks_list, output_chunks_dir, chunk_source_name
        )
        print_and_log(
            f"Chunked: split into this many chunks-> {number_of_chunks}",
            this_txt_output_dir_path,
        )

        append_chunks_to_jsonl(chunks_list, output_chunks_jsonl_path, chunk_source_name)

        print_and_log(f"{text_file_path} -> ok!", this_txt_output_dir_path)

        print("OK!")

    else:
        print_and_log(
            f"{text_file_path} -> Faile, no text extracted", this_txt_output_dir_path
        )


def extract_text_from_pdf(
    text_file_path,
    this_txt_output_dir_path,
    output_jsonl_path,
    output_json_dir,
    output_whole_txt_path,
    output_txt_dir,
    output_chunks_jsonl_path,
    output_chunks_dir,
    overlap_size=150,
    max_chunk_size=MAX_CHUNK_SIZE,
):
    """
    Extracts the text from a Microsoft Word (DOCX) file.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        str: The extracted text from the DOCX file.
    """

    ###################
    # Make Directories
    ###################

    # Create a directory for individual JSON files
    if not os.path.exists(output_json_dir):
        os.makedirs(output_json_dir)

    # Create a directory for individual txt files
    if not os.path.exists(output_txt_dir):
        os.makedirs(output_txt_dir)

    # Create a directory for chunks output_chunks_dir
    if not os.path.exists(output_chunks_dir):
        os.makedirs(output_chunks_dir)

    #########################
    # extract text from docx
    #########################
    text = simple_extracttextfrom_pdf(text_file_path)

    if text:
        # Save individual txt files
        text_file_path = os.path.basename(text_file_path)

        #################
        # .json & .jsonl
        #################

        # Write/Append to a single JSONL file
        with open(output_jsonl_path, "a") as f:
            json_record = json.dumps({"text": text.strip()})
            f.write(json_record + "\n")

        # Save individual JSON file
        individual_json_path = os.path.join(
            output_json_dir, f"{os.path.splitext(text_file_path)[0]}.json"
        )
        with open(individual_json_path, "w") as f:
            json.dump({"text": text.strip()}, f, indent=4)

        #######
        # .txt
        #######

        # Write/Append to a single text .txt file
        with open(output_whole_txt_path, "a") as f:
            f.write(text + "\n\n")

        # # Save individual txt files
        # individual_txt_path = os.path.join(output_txt_dir, f"{os.path.splitext(text_file_path)[0]}.txt")
        # print(f"individual_txt_path -> {individual_txt_path}")
        # with open(individual_txt_path, 'w') as f:
        #     f.write(text)

        # Save individual txt files
        base_name = os.path.basename(text_file_path)
        individual_txt_path = os.path.join(output_txt_dir, base_name)
        print(f"individual_txt_path -> {individual_txt_path}")
        with open(individual_txt_path, "w", encoding="utf-8") as f:
            f.write(text)

        #########
        # chunks
        #########

        chunks_list = make_chunk_list(
            text, max_chunk_size, this_txt_output_dir_path, overlap_size
        )

        chunk_source_name = os.path.splitext(text_file_path)[0]

        # check sizes
        check_len_chunks_in_list(chunks_list, max_chunk_size, this_txt_output_dir_path)

        number_of_chunks = save_individual_chunks(
            chunks_list, output_chunks_dir, chunk_source_name
        )
        print_and_log(
            f"Chunked: split into this many chunks-> {number_of_chunks}",
            this_txt_output_dir_path,
        )

        append_chunks_to_jsonl(chunks_list, output_chunks_jsonl_path, chunk_source_name)

        print_and_log(f"{text_file_path} -> ok!", this_txt_output_dir_path)

        print("OK!")

    else:
        print_and_log(
            f"{text_file_path} -> Faile, no text extracted", this_txt_output_dir_path
        )


def zip_folder(path_to_directory_to_zip, output_destination_zip_file_path):
    """Creates a zip archive of a specified folder.

    Args:
        path_to_directory_to_zip (str): The path to the folder to be zipped.
        output_destination_zip_file_path (str): The desired name and path of the output zip file.
    """
    # # Specify the folder you want to zip
    # path_to_directory_to_zip = "individual_jsons"

    # # Specify the desired output zip file name (e.g., 'jsons_archive.zip')
    # output_destination_zip_file_path = "jsons_archive_zip"

    shutil.make_archive(
        output_destination_zip_file_path, "zip", path_to_directory_to_zip
    )


###########
# Chunking
############


def split_sentences_and_punctuation(text):
    """Splits text into sentences, attempting to preserve punctuation and all text content.
    Args:
        text (str): The input text.
    Returns:
        list: A list of sentences with preserved punctuation.
    """
    ABBREVIATIONS = [
        "Dr.",
        "Mr.",
        "Mrs.",
        "Ms.",
        "Lt.",
        "St.",
        "Capt.",
        "Col.",
        "Gen.",
        "Rev.",
        "Hon.",
    ]

    # Construct a pattern to match abbreviations
    abbreviations_pattern = r"|".join(
        r"\b{}\b".format(re.escape(abbr)) for abbr in ABBREVIATIONS
    )

    # This pattern attempts to split at sentence endings (.?!), including the punctuation with the preceding sentence
    # It uses a lookahead to keep the punctuation with the sentence
    # The negative lookahead (?!({abbreviations_pattern})\s) excludes known abbreviations from being split
    sentence_end_regex = (
        r"(?<=[.!?])\s+(?=[A-Z])(?!({abbreviations_pattern})\s)".format(
            abbreviations_pattern=abbreviations_pattern
        )
    )

    split_sentences_and_punctuation_list = re.split(sentence_end_regex, text)

    # Optionally, remove empty strings if they are not desired
    split_sentences_and_punctuation_list = [
        s for s in split_sentences_and_punctuation_list if s
    ]

    return split_sentences_and_punctuation_list


# def split_sentences_and_punctuation(text):
#     """Splits text into sentences, attempting to preserve punctuation and all text content.

#     Args:
#         text (str): The input text.

#     Returns:
#         list: A list of sentences with preserved punctuation.
#     """

#     # This pattern attempts to split at sentence endings (.?!), including the punctuation with the preceding sentence
#     # It uses a lookahead to keep the punctuation with the sentence
#     sentence_end_regex = r'(?<=[.!?])\s+(?=[A-Z])'

#     split_sentences_and_punctuation_list = re.split(sentence_end_regex, text)

#     # Optionally, remove empty strings if they are not desired
#     split_sentences_and_punctuation_list = [s for s in split_sentences_and_punctuation_list if s]

#     # print("split_sentences_and_punctuation_list")
#     # print(split_sentences_and_punctuation_list)

#     return split_sentences_and_punctuation_list


def recombine_punctuation(sentences):
    """
    A helper function a that
    Recombines floating punctuation and
    creates a new list of sentences."""
    recombined_sentences = []
    i = 0

    while i < len(sentences) - 1:
        # print(i)
        # print(sentences[i-1])
        # print(sentences[i])
        # print(sentences[i+1])

        sentence = sentences[i].strip()
        next_item = sentences[i + 1].strip()

        # if next_item in ".?!":
        if re.match(r"[.?!]+", next_item):
            recombined_sentences.append(sentence + next_item)
            i += 1  # Skip the punctuation since it's been combined
        else:
            recombined_sentences.append(sentence)
        i += 1

    # Add the last sentence (if it exists)
    if sentences[-1]:
        recombined_sentences.append(sentences[-1].strip())

    return recombined_sentences


# # original, works with no overlap
# def chunk_text(sentences, chunk_size, overlap_size=0):
#     """
#     todo: add overlap
#     maybe by reserving the last sentence
#     """
#     chunked_text = []
#     current_chunk = ""
#     overlap_text = ""

#     for sentence in sentences:
#         # Case 1: Chunk + sentence easily fit
#         if len(current_chunk) + len(sentence) + 1 <= chunk_size:
#             current_chunk += sentence + " "

#         # Case 2: Sentence itself is too big
#         elif len(sentence) > chunk_size:
#             # Split long sentence (implement 'split_long_sentence' below)
#             for sub_sentence in split_long_sentence(sentence, chunk_size):
#                 chunked_text.append(sub_sentence.strip())

#         # Case 3:  Chunk + sentence exceed limit, time to split
#         else:
#             chunked_text.append(current_chunk.strip())
#             current_chunk = sentence + " "

#     # Handle final chunk
#     if current_chunk:
#         chunked_text.append(current_chunk.strip())

#     return chunked_text


def chunk_text(sentences, chunk_size, overlap_size=1000):
    """
    with overlap
    """
    overlap_size = 550

    chunk_size = chunk_size - overlap_size
    chunked_text = []
    current_chunk = ""
    last_sentence = ""

    for this_sentence in sentences:

        # Case 1: Chunk + sentence easily fit
        if len(current_chunk) + len(this_sentence) + 1 <= chunk_size:
            current_chunk += this_sentence + " "

        # Case 2: this_sentence itself is too big
        elif len(this_sentence) > chunk_size:
            # Split long sentence (implement 'split_long_sentence' below)
            for sub_sentence in split_long_sentence(this_sentence, chunk_size):
                chunked_text.append(sub_sentence.strip())
            current_chunk = ""

        # Case 3:  Chunk + this_sentence exceed limit, time to split
        else:
            chunked_text.append(current_chunk.strip())

            # start the next chunk (with the last_sentence)
            current_chunk = last_sentence + " " + this_sentence + " "

        last_sentence = this_sentence
        if len(this_sentence) > overlap_size:
            # just the last part
            last_sentence = this_sentence[-overlap_size:]

    # Handle final chunk
    if current_chunk:
        chunked_text.append(current_chunk.strip())

    return chunked_text


def split_long_sentence(sentence, chunk_size):
    """Splits a long sentence into chunks, aiming near the  chunk_size."""
    words = sentence.split()
    chunks = []
    current_chunk = ""

    for word in words:
        if len(current_chunk) + len(word) + 1 <= chunk_size:
            current_chunk += word + " "
        else:
            chunks.append(current_chunk.strip())
            current_chunk = word + " "

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks


def check_for_not(chunk, window_size=25):
    """Checks if 'not' is isolated near a potential cut.

    Args:
        chunk (list): A list of sentences forming the chunk.
        window_size (int): The number of characters to consider on either side.

    Returns:
        bool: True if 'not' is isolated, False otherwise.
    """

    joined_chunk = " ".join(chunk)
    not_indices = [m.start() for m in re.finditer(r"\bnot\b", joined_chunk)]

    for index in not_indices:
        start = max(0, index - window_size)
        end = min(len(joined_chunk), index + window_size)
        if not re.search(r"\w", joined_chunk[start:end]):  # Check for surrounding words
            return True

    return False


def make_chunk_list(text, chunk_size, this_epub_output_dir_path, overlap_size=150):
    split_sentences_list = split_sentences_and_punctuation(text)
    chunk_list = chunk_text(split_sentences_list, chunk_size, overlap_size)

    for i in chunk_list:
        if not i:
            print_and_log(
                "error None in chunk_list: make_chunk_list()", this_epub_output_dir_path
            )

    print_and_log(f"len chunk list -> {len(chunk_list)}", this_epub_output_dir_path)

    return chunk_list


#######
# pool
#######


# def pool_txt_files(src_dir, dest_dir):
#     """
#     Recursively copies all .txt files from the source directory to the destination directory.
#     """
#     counter = 0

#     if not os.path.exists(dest_dir):
#         try:
#             os.makedirs(dest_dir)
#         except OSError as e:
#             if e.errno != errno.EEXIST:
#                 raise

#     for root, dirs, files in os.walk(src_dir):
#         for file in files:
#             if file.endswith('.txt'):
#                 src_file = os.path.join(root, file)
#                 dest_file = os.path.join(dest_dir, os.path.relpath(src_file, src_dir))
#                 os.makedirs(os.path.dirname(dest_file), exist_ok=True)
#                 shutil.copy2(src_file, dest_file)
#                 # print(f"Copied {src_file} to {dest_file}")

#                 counter += 1


#     print(f"\nOK! I just copied {counter} files, to {dest_dir}.")


def add_source_attribution(source_attribution_string, directory="txt_pool"):
    """
    This function adds a source attribution string to the end of each file in a directory.
    The attribution string is extracted from the filename.

    Parameters:
    directory (str): The path to the directory containing the files.
    """

    # Loop through each file in the directory
    for filename in os.listdir(directory):

        # Check if the item is a file (not a directory)
        if os.path.isfile(os.path.join(directory, filename)):

            # Extract the number from the filename
            match = re.search(r"_(\d+)\.txt$", filename)
            if match:
                number = match.group(1)

                # Create the attribution string to add to the file
                file_source_attribution_string = (
                    f"\n(Source: {source_attribution_string} Excerpt number {number})\n"
                )

                # Open the file for reading and writing
                with open(os.path.join(directory, filename), "a") as file:

                    # print(file_source_attribution_string)

                    # Add the attribution string to the end of the file
                    file.write(file_source_attribution_string)


def back_append_stub_files(directory):
    """
    This function processes all the .txt files in a given directory.
    If a file is smaller than 300 bytes, its content is appended to the previous file,
    and the current file is deleted. The process is then restarted from the beginning of the list.

    :param directory: The directory containing the .txt files
    """
    try:
        print(f"starting back_append_stub_files(directory)... directory -> {directory}")

        counter = 0

        # Get a list of all the .txt files in the directory and sort them
        files = sorted([f for f in os.listdir(directory) if f.endswith(".txt")])

        i = 1
        while i < len(files):
            # Get the full path of the current file
            current_file = os.path.join(directory, files[i])

            # Get the size of the current file
            size = os.path.getsize(current_file)

            # If the size of the current file is smaller than 300 bytes
            if size < MINIMUM_BYTES_SIZE:
                # Read the content of the current file
                with open(current_file, "r") as f:
                    content = f.read()

                # Get the full path of the previous file
                previous_file = os.path.join(directory, files[i - 1])

                # Append the content of the current file to the previous file
                with open(previous_file, "a") as f:
                    f.write("\n" + content)

                # Delete the current file
                os.remove(current_file)

                # Remove the current file from the list
                files.pop(i)

                counter += 1
            else:
                # Move on to the next file
                i += 1

        print(f"Back appended {counter} files.")

    except Exception as e:
        raise e


def pool_txt_files(src_dir, dest_dir):
    """
    Recursively copies all .txt files
    from the source directory to the destination directory.
    original files are not deleted.
    """
    counter = 0

    # Create the destination directory if it doesn't exist
    if not os.path.exists(dest_dir):
        os.makedirs(dest_dir)

    # Traverse the source directory recursively
    for root, _, files in os.walk(src_dir):
        for file in files:
            if file.endswith(".txt"):
                src_file = os.path.join(root, file)
                dest_file = os.path.join(dest_dir, file)
                shutil.copy2(src_file, dest_file)
                counter += 1

    print(f"\nOK! I just copied {counter} files to {dest_dir}.")


def rename_directory(current_name, new_name):
    """
    # Usage
    rename_directory('old_name', 'new_name')
    """
    try:
        os.rename(current_name, new_name)
        print(f"Directory renamed from {current_name} to {new_name}")
    except FileNotFoundError:
        print(f"Directory {current_name} not found")
    except Exception as e:
        print(f"An error occurred: {e}")


def make_merged_directories_pool():
    # Create the 'txt_pool' directory if it doesn't exist
    if not os.path.exists("txt_pool"):
        os.makedirs("txt_pool")

    # Iterate over all directories that start with 'txt_pool_'
    for dir_name in os.listdir("."):
        if dir_name.startswith("txt_pool_"):
            # Iterate over all files in the directory
            for file_name in os.listdir(dir_name):
                # Move the file to the 'txt_pool' directory
                shutil.copy(os.path.join(dir_name, file_name), "txt_pool")


######
# Run
######
"""
1. add your epub files into the same current working directory as this script
2. run script
3. find the files in new folders per epub
"""


################################
# Set source_attribution_string
################################


# Get the current working directory.
cwd = os.getcwd()

# # Search for all x-type files in the current working directory.
# epub_files = glob.glob(os.path.join(cwd, "*.epub"))

base_dir = "target_files"
file_extension = ".epub"

epub_files = get_files_from_subdirs(base_dir, file_extension)

# Print the list of EPUB files.
print("epub files")
print(epub_files)


####################
# run for each epub
####################
for this_epub_file in epub_files:
    print(f"\n\n For epub: {this_epub_file}")

    source_attribution_string = ""

    source_attribution_string = input(
        "What is your source attribution string? (e.g. title, author, publisher, year, etc.)\n (enter to use file name)\n"
    )

    # use file path by default
    if source_attribution_string == "":
        file_name = os.path.basename(this_epub_file)
        source_attribution_string = file_name

    # set target epub to first epub doc listed as being in the cwd
    epub_file_path = this_epub_file

    # make directory for this book
    this_epub_output_dir_path = epub_file_path[:-5] + "_epub_folder"
    print(this_epub_output_dir_path)

    # Add another parent directory
    parent_dir = RESULTS_DIR_NAME
    this_epub_output_dir_path = os.path.join(parent_dir, this_epub_output_dir_path)

    # Set the absolute path
    this_epub_output_dir_path = os.path.abspath(this_epub_output_dir_path)

    # Create a directory for individual txt files
    if not os.path.exists(this_epub_output_dir_path):
        os.makedirs(this_epub_output_dir_path)

    # json
    # output_jsonl_path = 'output.jsonl'
    output_jsonl_path = os.path.join(this_epub_output_dir_path, "output.jsonl")
    output_json_dir = os.path.join(
        this_epub_output_dir_path, "individual_jsons"
    )  # Directory to store individual JSON files
    output_json_zip_dir = os.path.join(
        this_epub_output_dir_path, "jsons_zip_archive"
    )  # Directory to store individual JSON files

    # txt
    output_whole_txt_path = os.path.join(this_epub_output_dir_path, "whole.txt")
    output_txt_dir = os.path.join(
        this_epub_output_dir_path, "individual_txt"
    )  # Directory to store individual txt files
    output_txt_zip_dir = os.path.join(
        this_epub_output_dir_path, "txt_zip_archive"
    )  # Directory to store individual JSON files

    # chunks
    output_chunks_jsonl_path = os.path.join(
        this_epub_output_dir_path, "chunks_jsonl_all.jsonl"
    )  # Directory to store individual txt files
    output_chunks_dir = os.path.join(
        this_epub_output_dir_path, "chunk_text_files"
    )  # Directory to store individual txt files
    output_chunks_zip_dir = os.path.join(
        this_epub_output_dir_path, "chunks_zip_archive"
    )  # Directory to store individual JSON files

    extract_text_from_epub(
        epub_file_path,
        this_epub_output_dir_path,
        output_jsonl_path,
        output_json_dir,
        output_whole_txt_path,
        output_txt_dir,
        output_chunks_jsonl_path,
        output_chunks_dir,
        max_chunk_size=MAX_CHUNK_SIZE,
    )

    # Call the zip function
    """
    zip_folder(path_to_directory_to_zip, output_destination_zip_file_path)
    """
    zip_folder(output_json_dir, output_json_zip_dir)
    zip_folder(output_txt_dir, output_txt_zip_dir)
    zip_folder(output_chunks_dir, output_chunks_zip_dir)

    ######################################
    # Bundle of Additional/Optional Items
    ######################################

    # Removes files smaller than N bytes from X directory
    remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

    # Call the function with the directory containing the .txt files
    back_append_stub_files("txt_pool")

    add_source_attribution(source_attribution_string)

    # count all file chunks created
    count_files("txt_pool")

    # rename_directory('old_name', 'new_name')
    rename_directory("txt_pool", f"txt_pool_{pool_counter}")

    pool_counter += 1


######
# Run txt
######
"""
1. add your .txt files into the same current working directory as this script
2. run script
3. find the files in new folders per epub
"""


# Get the current working directory.
cwd = os.getcwd()

# # Search for all x-type files in the current working directory.
# txt_files = glob.glob(os.path.join(cwd, "*.txt"))

base_dir = "target_files"
file_extension = ".txt"

# recoursive file search with THE GLOB
txt_files = get_files_from_subdirs(base_dir, file_extension)

# Print the list of txt files.
print("text files")
print(txt_files)


####################
# run for each txt
####################
for this_txt_file in txt_files:
    print(f"\n\n For txt: {this_txt_file}")

    source_attribution_string = ""

    source_attribution_string = input(
        "What is your source attribution string? (e.g. title, author, publisher, year, etc.)\n (enter to use file name)\n"
    )

    # use file path by default
    if source_attribution_string == "":
        file_name = os.path.basename(this_txt_file)
        source_attribution_string = file_name

    # set target txt to first txt doc listed as being in the cwd
    text_file_path = this_txt_file

    # make directory for this book
    this_txt_output_dir_path = text_file_path[:-5] + "_txt_folder"
    print(this_txt_output_dir_path)

    # Add another parent directory
    parent_dir = RESULTS_DIR_NAME
    this_txt_output_dir_path = os.path.join(parent_dir, this_txt_output_dir_path)

    # Set the absolute path
    this_txt_output_dir_path = os.path.abspath(this_txt_output_dir_path)

    # Create a directory for individual txt files
    if not os.path.exists(this_txt_output_dir_path):
        os.makedirs(this_txt_output_dir_path)

    # json
    # output_jsonl_path = 'output.jsonl'
    output_jsonl_path = os.path.join(this_txt_output_dir_path, "output.jsonl")
    output_json_dir = os.path.join(
        this_txt_output_dir_path, "individual_jsons"
    )  # Directory to store individual JSON files
    output_json_zip_dir = os.path.join(
        this_txt_output_dir_path, "jsons_zip_archive"
    )  # Directory to store individual JSON files

    # txt
    output_whole_txt_path = os.path.join(this_txt_output_dir_path, "whole.txt")
    output_txt_dir = os.path.join(
        this_txt_output_dir_path, "individual_txt"
    )  # Directory to store individual txt files
    output_txt_zip_dir = os.path.join(
        this_txt_output_dir_path, "txt_zip_archive"
    )  # Directory to store individual JSON files

    # chunks
    output_chunks_jsonl_path = os.path.join(
        this_txt_output_dir_path, "chunks_jsonl_all.jsonl"
    )  # Directory to store individual txt files
    output_chunks_dir = os.path.join(
        this_txt_output_dir_path, "chunk_text_files"
    )  # Directory to store individual txt files
    output_chunks_zip_dir = os.path.join(
        this_txt_output_dir_path, "chunks_zip_archive"
    )  # Directory to store individual JSON files

    extract_text_from_txt(
        text_file_path,
        this_txt_output_dir_path,
        output_jsonl_path,
        output_json_dir,
        output_whole_txt_path,
        output_txt_dir,
        output_chunks_jsonl_path,
        output_chunks_dir,
        max_chunk_size=MAX_CHUNK_SIZE,
    )

    # Call the zip function
    """
    zip_folder(path_to_directory_to_zip, output_destination_zip_file_path)
    """
    zip_folder(output_json_dir, output_json_zip_dir)
    zip_folder(output_txt_dir, output_txt_zip_dir)
    zip_folder(output_chunks_dir, output_chunks_zip_dir)

    ######################################
    # Bundle of Additional/Optional Items
    ######################################

    # Removes files smaller than N bytes from X directory
    remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

    # Call the function with the directory containing the .txt files
    back_append_stub_files("txt_pool")

    add_source_attribution(source_attribution_string)

    # count all file chunks created
    count_files("txt_pool")

    # rename_directory('old_name', 'new_name')
    rename_directory("txt_pool", f"txt_pool_{pool_counter}")

    pool_counter += 1


######
# Run docx
######
"""
1. add your .txt files into the same current working directory as this script
2. run script
3. find the files in new folders per epub
"""


# Get the current working directory.
cwd = os.getcwd()

# # Search for all x-type files in the current working directory.
# docx_files = glob.glob(os.path.join(cwd, "*.docx"))

base_dir = "target_files"
file_extension = ".docx"

# recoursive file search with THE GLOB
docx_files = get_files_from_subdirs(base_dir, file_extension)

# Print the list of docx files.
print("docx_files")
print(docx_files)


####################
# run for each docx
####################
for this_docx_file in docx_files:
    print(f"\n\n For docx: {this_docx_file}")

    source_attribution_string = ""

    source_attribution_string = input(
        "What is your source attribution string? (e.g. title, author, publisher, year, etc.)\n (enter to use file name)\n"
    )

    # use file path by default
    if source_attribution_string == "":
        file_name = os.path.basename(this_docx_file)
        source_attribution_string = file_name

    # set target txt to first txt doc listed as being in the cwd
    text_file_path = this_docx_file

    # make directory for this book
    this_txt_output_dir_path = text_file_path[:-5] + "_docx_folder"
    print(this_txt_output_dir_path)

    # Add another parent directory
    parent_dir = RESULTS_DIR_NAME
    this_txt_output_dir_path = os.path.join(parent_dir, this_txt_output_dir_path)

    # Set the absolute path
    this_txt_output_dir_path = os.path.abspath(this_txt_output_dir_path)

    # Create a directory for individual txt files
    if not os.path.exists(this_txt_output_dir_path):
        os.makedirs(this_txt_output_dir_path)

    # json
    # output_jsonl_path = 'output.jsonl'
    output_jsonl_path = os.path.join(this_txt_output_dir_path, "output.jsonl")
    output_json_dir = os.path.join(
        this_txt_output_dir_path, "individual_jsons"
    )  # Directory to store individual JSON files
    output_json_zip_dir = os.path.join(
        this_txt_output_dir_path, "jsons_zip_archive"
    )  # Directory to store individual JSON files

    # txt
    output_whole_txt_path = os.path.join(this_txt_output_dir_path, "whole.txt")
    output_txt_dir = os.path.join(
        this_txt_output_dir_path, "individual_txt"
    )  # Directory to store individual txt files
    output_txt_zip_dir = os.path.join(
        this_txt_output_dir_path, "txt_zip_archive"
    )  # Directory to store individual JSON files

    # chunks
    output_chunks_jsonl_path = os.path.join(
        this_txt_output_dir_path, "chunks_jsonl_all.jsonl"
    )  # Directory to store individual txt files
    output_chunks_dir = os.path.join(
        this_txt_output_dir_path, "chunk_text_files"
    )  # Directory to store individual txt files
    output_chunks_zip_dir = os.path.join(
        this_txt_output_dir_path, "chunks_zip_archive"
    )  # Directory to store individual JSON files

    extract_text_from_docx(
        text_file_path,
        this_txt_output_dir_path,
        output_jsonl_path,
        output_json_dir,
        output_whole_txt_path,
        output_txt_dir,
        output_chunks_jsonl_path,
        output_chunks_dir,
        max_chunk_size=MAX_CHUNK_SIZE,
    )

    # Call the zip function
    """
    zip_folder(path_to_directory_to_zip, output_destination_zip_file_path)
    """
    zip_folder(output_json_dir, output_json_zip_dir)
    zip_folder(output_txt_dir, output_txt_zip_dir)
    zip_folder(output_chunks_dir, output_chunks_zip_dir)

    ######################################
    # Bundle of Additional/Optional Items
    ######################################

    # Removes files smaller than N bytes from X directory
    remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

    # Call the function with the directory containing the .txt files
    back_append_stub_files("txt_pool")

    add_source_attribution(source_attribution_string)

    # count all file chunks created
    count_files("txt_pool")

    # rename_directory('old_name', 'new_name')
    rename_directory("txt_pool", f"txt_pool_{pool_counter}")

    pool_counter += 1

######
# Run pdf
######
"""
1. add your .txt files into the same current working directory as this script
2. run script
3. find the files in new folders per epub
"""

overlap_size = 300

# Get the current working directory.
cwd = os.getcwd()

# Search for all x-type files in the current working directory.
pdf_files = glob.glob(os.path.join(cwd, "*.pdf")) + glob.glob(
    os.path.join(cwd, "*.PDF")
)

base_dir = "target_files"
file_extension = ".pdf"
# recoursive file search with THE GLOB
pdf_files1 = get_files_from_subdirs(base_dir, file_extension)


base_dir = "target_files"
file_extension = ".PDF"
# recoursive file search with THE GLOB
pdf_files2 = get_files_from_subdirs(base_dir, file_extension)

# both alt versions of suffix
pdf_files = pdf_files1 + pdf_files2

set_pdf_files = set(pdf_files)
pdf_files = list(set_pdf_files)

# Print the list of txt files.
print("pdf_files")
print(pdf_files)


####################
# run for each pdf
####################
for this_pdf_file in pdf_files:
    print(f"\n\n For pdf_files: {this_pdf_file}")

    source_attribution_string = ""

    source_attribution_string = input(
        "What is your source attribution string? (e.g. title, author, publisher, year, etc.)\n (enter to use file name)\n"
    )

    # use file path by default
    if source_attribution_string == "":
        file_name = os.path.basename(this_pdf_file)
        source_attribution_string = file_name

    # set target txt to first txt doc listed as being in the cwd
    text_file_path = this_pdf_file

    # make directory for this book
    this_txt_output_dir_path = text_file_path[:-5] + "_pdf_folder"
    print(this_txt_output_dir_path)

    # Add another parent directory
    parent_dir = RESULTS_DIR_NAME
    this_txt_output_dir_path = os.path.join(parent_dir, this_txt_output_dir_path)

    # Set the absolute path
    this_txt_output_dir_path = os.path.abspath(this_txt_output_dir_path)

    # Create a directory for individual txt files
    if not os.path.exists(this_txt_output_dir_path):
        os.makedirs(this_txt_output_dir_path)

    # json
    # output_jsonl_path = 'output.jsonl'
    output_jsonl_path = os.path.join(this_txt_output_dir_path, "output.jsonl")
    output_json_dir = os.path.join(
        this_txt_output_dir_path, "individual_jsons"
    )  # Directory to store individual JSON files
    output_json_zip_dir = os.path.join(
        this_txt_output_dir_path, "jsons_zip_archive"
    )  # Directory to store individual JSON files

    # txt
    output_whole_txt_path = os.path.join(this_txt_output_dir_path, "whole.txt")
    output_txt_dir = os.path.join(
        this_txt_output_dir_path, "individual_txt"
    )  # Directory to store individual txt files
    output_txt_zip_dir = os.path.join(
        this_txt_output_dir_path, "txt_zip_archive"
    )  # Directory to store individual JSON files

    # chunks
    output_chunks_jsonl_path = os.path.join(
        this_txt_output_dir_path, "chunks_jsonl_all.jsonl"
    )  # Directory to store individual txt files
    output_chunks_dir = os.path.join(
        this_txt_output_dir_path, "chunk_text_files"
    )  # Directory to store individual txt files
    output_chunks_zip_dir = os.path.join(
        this_txt_output_dir_path, "chunks_zip_archive"
    )  # Directory to store individual JSON files

    extract_text_from_pdf(
        text_file_path,
        this_txt_output_dir_path,
        output_jsonl_path,
        output_json_dir,
        output_whole_txt_path,
        output_txt_dir,
        output_chunks_jsonl_path,
        output_chunks_dir,
        overlap_size,
        max_chunk_size=MAX_CHUNK_SIZE,
    )

    # Call the zip function
    """
    zip_folder(path_to_directory_to_zip, output_destination_zip_file_path)
    """
    zip_folder(output_json_dir, output_json_zip_dir)
    zip_folder(output_txt_dir, output_txt_zip_dir)
    zip_folder(output_chunks_dir, output_chunks_zip_dir)

    ######################################
    # Bundle of Additional/Optional Items
    ######################################

    # Removes files smaller than N bytes from X directory
    remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

    # Call the function with the directory containing the .txt files
    back_append_stub_files("txt_pool")

    add_source_attribution(source_attribution_string)

    # count all file chunks created
    count_files("txt_pool")

    # rename_directory('old_name', 'new_name')
    rename_directory("txt_pool", f"txt_pool_{pool_counter}")

    pool_counter += 1


#######
# pptx
#######
# directory_to_process = "files_to_process"
# pptx_files = glob.glob(f"{directory_to_process}/*.pptx")

pptx_base_dir = "target_files"
file_extension = ".pptx"
# recoursive file search with THE GLOB
pptx_files = get_files_from_subdirs(pptx_base_dir, file_extension)

print("pptx_files")
print(pptx_files)

# Process each .pptx file
for file_path in pptx_files:
    print(f"\n\n For pptx_files: {file_path}")

    source_attribution_string = ""

    source_attribution_string = input(
        "What is your source attribution string? (e.g. title, author, publisher, year, etc.)\n (enter to use file name)\n"
    )

    # use file path by default
    if source_attribution_string == "":
        file_name = os.path.basename(file_path)
        source_attribution_string = file_name

    extract_pptx_text_to_file(file_path)

    ######################################
    # Bundle of Additional/Optional Items
    ######################################

    # Removes files smaller than N bytes from X directory
    remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

    # Call the function with the directory containing the .txt files
    back_append_stub_files("txt_pool")

    add_source_attribution(source_attribution_string)

    # count all file chunks created
    count_files("txt_pool")

    # rename_directory('old_name', 'new_name')
    rename_directory("txt_pool", f"txt_pool_{pool_counter}")

    pool_counter += 1


# move pptx files
pool_txt_files("pptx_chunks", "txt_pool")

# # make sure all files in pool are not nested
# pool_txt_files(txt_pool, 'txt_pool_all')


######################################
# Bundle of Additional/Optional Items
######################################

# Removes files smaller than N bytes from X directory
remove_small_files("txt_pool", REMOVE_BELOW_SIZE)

# Call the function with the directory containing the .txt files
back_append_stub_files("txt_pool")

add_source_attribution(source_attribution_string)

# count all file chunks created
count_files("txt_pool")

# rename_directory('old_name', 'new_name')
rename_directory("txt_pool", f"txt_pool_{pool_counter}")

pool_counter += 1


##############
# merged pool
##############
make_merged_directories_pool()


# end & print debug timer
end_time_inner = time.time()
end_time_outer = time.monotonic()
elapsed_time = end_time_inner - start_time_inner
print(
    f"Smart-Chunk automated document processing: Inner Elapsed time: {elapsed_time} seconds"
)
elapsed_time = end_time_outer - start_time_outer
print(
    f"Smart-Chunk automated document processing: Outer Elapsed time: {elapsed_time} seconds"
)

# count all file chunks created
count_files("txt_pool")
